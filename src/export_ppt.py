from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PURPLE = RGBColor(0xA1, 0x00, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_PURPLE = RGBColor(0xE8, 0xC0, 0xFF)


def _add_title_slide(prs: Presentation, client_name: str) -> None:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    # Purple rectangle as accent (simulating triangle)
    triangle = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE freeform — use rectangle
        Inches(0.4), Inches(0.3), Inches(0.8), Inches(0.8)
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = PURPLE
    triangle.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.4), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Tech Debt Assessment"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Client name
    txBox2 = slide.shapes.add_textbox(Inches(1.4), Inches(1.2), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = client_name or "Client Assessment"
    p2.font.size = Pt(18)
    p2.font.color.rgb = PURPLE

    # Purple line
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(2.0), Inches(9.2), Emu(40000))
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()


def _add_kpi_slide(prs: Presentation, results: dict) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Summary"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    scores = results.get("scores", {})
    kpis = [
        ("Application Score", f"{scores.get('application', 0):.2f}"),
        ("Infrastructure Score", f"{scores.get('infrastructure', 0):.2f}"),
        ("Architecture Score", f"{scores.get('architecture', 0):.2f}"),
        ("People Score", f"{scores.get('people', 0):.2f}"),
        ("Total Annual Interest", f"{results.get('total_interest', 0):.1f} kUSD"),
        ("Total TCO", f"{results.get('total_tco', 0):.1f} kUSD/year"),
    ]

    cols = 3
    box_w = Inches(2.8)
    box_h = Inches(1.4)
    start_x = Inches(0.4)
    start_y = Inches(1.0)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for i, (label, value) in enumerate(kpis):
        col = i % cols
        row = i // cols
        x = start_x + col * (box_w + gap_x)
        y = start_y + row * (box_h + gap_y)

        box = slide.shapes.add_shape(1, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_GREY
        box.line.color.rgb = PURPLE

        txBox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), box_w - Inches(0.2), box_h - Inches(0.2))
        tf = txBox.text_frame
        tf.word_wrap = True

        p_label = tf.paragraphs[0]
        p_label.text = label
        p_label.font.size = Pt(10)
        p_label.font.color.rgb = LIGHT_PURPLE

        p_val = tf.add_paragraph()
        p_val.text = value
        p_val.font.size = Pt(18)
        p_val.font.bold = True
        p_val.font.color.rgb = WHITE


def _add_interest_slide(prs: Presentation, results: dict) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Annual Interest Cost by Dimension"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    interest_dim = results.get("interest_by_dimension", {})
    dims = [
        ("Application", interest_dim.get("application", 0)),
        ("Infrastructure", interest_dim.get("infrastructure", 0)),
        ("Architecture", interest_dim.get("architecture", 0)),
        ("People", interest_dim.get("people", 0)),
    ]
    total = results.get("total_interest", 0)
    max_val = max(v for _, v in dims) if any(v for _, v in dims) else 1

    bar_start_x = Inches(1.5)
    bar_max_w = Inches(7.0)
    bar_h = Inches(0.5)
    gap_y = Inches(0.3)
    start_y = Inches(1.2)
    purple_variants = [PURPLE, RGBColor(0xCC, 0x00, 0xFF), RGBColor(0x7B, 0x00, 0xCC), RGBColor(0x55, 0x00, 0x88)]

    for i, (label, value) in enumerate(dims):
        y = start_y + i * (bar_h + gap_y)

        lbl_box = slide.shapes.add_textbox(Inches(0.4), y, Inches(1.0), bar_h)
        lf = lbl_box.text_frame
        lp = lf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(11)
        lp.font.color.rgb = WHITE

        bar_w = bar_max_w * (value / max_val) if max_val > 0 else Inches(0.1)
        bar = slide.shapes.add_shape(1, bar_start_x, y, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = purple_variants[i]
        bar.line.fill.background()

        val_box = slide.shapes.add_textbox(bar_start_x + bar_w + Inches(0.1), y, Inches(1.5), bar_h)
        vf = val_box.text_frame
        vp = vf.paragraphs[0]
        vp.text = f"{value:.1f} kUSD"
        vp.font.size = Pt(11)
        vp.font.color.rgb = LIGHT_PURPLE

    total_y = start_y + len(dims) * (bar_h + gap_y) + Inches(0.2)
    tot_box = slide.shapes.add_textbox(Inches(0.4), total_y, Inches(9), Inches(0.5))
    tf = tot_box.text_frame
    tp = tf.paragraphs[0]
    tp.text = f"Total Annual Interest: {total:.1f} kUSD"
    tp.font.size = Pt(14)
    tp.font.bold = True
    tp.font.color.rgb = PURPLE


def generate_ppt(assessment: dict, results: dict) -> bytes:
    """Return PPT bytes for the Accenture-branded presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    client_name = assessment.get("client", {}).get("name", "")
    _add_title_slide(prs, client_name)
    _add_kpi_slide(prs, results)
    _add_interest_slide(prs, results)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
